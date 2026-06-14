#!/usr/bin/env python
"""
根据 PARAMETERS.md 生成汇报 PPT —— 每个参数一页幻灯片。

运行方式:
    python tools/gen_pptx.py
    python tools/gen_pptx.py -o ./output/report.pptx
"""

import os, sys, argparse
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SHOWCASE_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                            "output", "param_showcase")
IMGDIR = SHOWCASE_DIR

# ── 样式 ──────────────────────────────────────────────
C_PRI   = RGBColor(0x1A, 0x56, 0xDB)
C_ACC   = RGBColor(0xE6, 0x8A, 0x2E)
C_DARK  = RGBColor(0x2D, 0x2D, 0x2D)
C_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED = RGBColor(0x88, 0x88, 0x88)
C_BG    = RGBColor(0xFA, 0xFA, 0xFA)

# ── 每个参数的解释文本 ──────────────────────────────
GEN_EXPLANATIONS = {
    "color_std": (
        "畴区间颜色差异标准差。值越大各畴区颜色差异越明显。\n"
        "baseline=0.2（固定），畴区间颜色几乎一致，模拟单晶材料。"
    ),
    "texture_std": (
        "畴区内部逐像素高斯纹理噪声标准差。值越大畴区表面越粗糙。\n"
        "baseline=2，产生轻微内部纹理，模拟天然晶粒表面。"
    ),
    "bg_noise_std": (
        "背景逐像素高斯噪声标准差，模拟传感器暗电流噪声或基底纹理。\n"
        "baseline=2，模拟典型显微镜背景噪声。"
    ),
    "color_mean": (
        "畴区颜色的 RGB 均值。每个畴区在此基础色上叠加 color_std 随机偏移。\n"
        "baseline=(163,138,127) 暖棕色。可根据目标材料调整。"
    ),
    "bg_mean": (
        "背景颜色的 RGB 均值。每个背景像素叠加 bg_noise_std 高斯噪声。\n"
        "baseline=(153,115,85)，比畴区略暗以形成区分度。"
    ),
    "base_r_range": (
        "未缩放时的畴区半径范围 (px)。最终半径 = base_r × mag × image_scale × ss。\n"
        "baseline=(15,25)。配合 mag_factor 和 image_scale 达到目标畴区尺寸。"
    ),
    "shape_jitter": (
        "六边形顶点径向扰动比例。0=正六边形，越大形状越不规则。\n"
        "baseline=0.1，模拟天然晶粒的轻微不规则性。"
    ),
    "orientation_std": (
        "畴区取向分布的标准差（度）。0=所有畴区同向，值越大取向越随机。\n"
        "baseline=0（同向）。5-15° 模拟轻微取向偏差，30-60° 模拟多晶随机取向。"
    ),
    "size_std": (
        "畴区半径分布的标准差。None=均匀分布，数值=正态分布。\n"
        "baseline=25（正态），模拟天然畴区的正态大小分布。"
    ),
    "base_num_range": (
        "未缩放时每张图的畴区数量范围。最终数量 = base_num / mag_factor²。\n"
        "baseline=(200,400)，在 10× 参考倍率下生效。实际数量受重叠控制器约束。"
    ),
    "edge_burr_amplitude": (
        "边缘毛刺振幅（相对边长），沿法向随机游走扰动模拟晶体生长粗糙度。\n"
        "baseline=0.08（20× 视场展示），产生自然的晶粒边缘。0=光滑边缘。"
    ),
    "edge_burr_subdivisions": (
        "每条边细分点数，控制毛刺的细密程度。≤1 时完全禁用。\n"
        "baseline=4（20× 视场展示），配合 amplitude=0.08 产生自然毛刺。"
    ),
    "max_overlap_ratio": (
        "新畴区与已有畴区交叠面积占自身面积的最大允许比例 [0,1]。\n"
        "baseline=0.3（20× 视场展示）。0=禁止重叠，1=允许完全重叠。"
    ),
    "max_overlap_count": (
        "任一像素最多允许被多少个畴区覆盖。None=不限制，1=禁止重叠。\n"
        "baseline=3（20× 视场展示）。控制畴区层叠深度。"
    ),
    "contain_threshold": (
        "包含判定阈值。新畴区覆盖某旧畴区超过该比例则拒绝放置，防止吞没。\n"
        "baseline=0.85（20× 视场展示）。0.7-0.9 为合理区间。"
    ),
    "supersample_ratio": (
        "超采样抗锯齿倍率。内部以 N× 尺寸渲染后 Lanczos 降采样。\n"
        "baseline=1（无抗锯齿）。N=2 时渲染像素为 4×，边缘锯齿明显时使用。"
    ),
}

AUG_EXPLANATIONS = {
    "Brightness": (
        "通过 TF.adjust_brightness() 调整图像亮度。因子 <1 变暗，>1 变亮。\n"
        "gen_multi_mag 默认: range=(0.8,1), prob=0.8（仅向暗调变化）。"
    ),
    "Contrast": (
        "通过 TF.adjust_contrast() 调整对比度。因子 <1 降低，>1 增强。\n"
        "gen_multi_mag 默认: range=(0.4,1), prob=0.8（仅向低对比度变化）。"
    ),
    "Gamma": (
        "通过 TF.adjust_gamma() 非线性亮度映射。<1 提亮暗部，>1 压暗暗部。\n"
        "gen_multi_mag 默认: range=(0.7,1.3), prob=0.8。"
    ),
    "Gaussian": (
        "通过 PIL.ImageFilter.GaussianBlur() 模拟显微镜失焦/光学衍射。\n"
        "gen_multi_mag 默认: range=(0.5,1), prob=1.0（始终触发）。"
    ),
    "Salt": (
        "随机将像素置为纯白(盐)或纯黑(胡椒)，模拟传感器坏点/脉冲噪声。\n"
        "gen_multi_mag 默认: prob=0.0（关闭）。amount 典型值 0.001-0.02。"
    ),
    "Color": (
        "在每个 RGB 通道施加全局偏移，模拟光源色温/白平衡偏差。\n"
        "在所有其他增强之前施加。gen_multi_mag 默认: range=(-12,12), prob=0.8。"
    ),
}

# 为每个生成器参数映射到图片文件名中的关键词
GEN_PARAM_IMG_MAP = {
    "color_std":             "gen_color_std_—_Inter-domain_Color_Variation.png",
    "texture_std":           "gen_texture_std_—_Intra-domain_Texture_Noise.png",
    "bg_noise_std":          "gen_bg_noise_std_—_Background_Noise.png",
    "color_mean":            "gen_color_mean_—_Domain_Hue.png",
    "bg_mean":               "gen_bg_mean_—_Background_Hue.png",
    "base_r_range":          "gen_base_r_range_—_Domain_Radius_Range.png",
    "shape_jitter":          "gen_shape_jitter_—_Vertex_Radial_Perturbation.png",
    "orientation_std":       "gen_orientation_std_—_Domain_Orientation_Spread.png",
    "size_std":              "gen_size_std_—_Size_Distribution.png",
    "base_num_range":        "gen_base_num_range_—_Domain_Count_Range.png",
    "edge_burr_amplitude":   "gen_edge_burr_amplitude_—_Edge_Roughness_Amplitude_20×.png",
    "edge_burr_subdivisions":"gen_edge_burr_subdivisions_—_Edge_Subdivision_Density_20×.png",
    "max_overlap_ratio":     "gen_max_overlap_ratio_—_Max_Overlap_Ratio_20×.png",
    "max_overlap_count":     "gen_max_overlap_count_—_Max_Pixel_Coverage_Layers_20×.png",
    "contain_threshold":     "gen_contain_threshold_—_Containment_Threshold_20×.png",
    "supersample_ratio":     "gen_supersample_ratio_—_Anti-aliasing_Supersampling.png",
}

AUG_PARAM_IMG_MAP = {
    "Brightness": "aug_Brightness_Adjustment.png",
    "Contrast":   "aug_Contrast_Adjustment.png",
    "Gamma":      "aug_Gamma_Correction.png",
    "Gaussian":   "aug_Gaussian_Blur.png",
    "Salt":       "aug_Salt_&_Pepper_Noise.png",
    "Color":      "aug_Color_Jitter___White_Balance_Shift.png",
}

AUG_PARAM_NAMES = {
    "Brightness": "Brightness — 亮度调整",
    "Contrast":   "Contrast — 对比度调整",
    "Gamma":      "Gamma — Gamma 校正",
    "Gaussian":   "Gaussian Blur — 高斯模糊",
    "Salt":       "Salt & Pepper Noise — 椒盐噪声",
    "Color":      "Color Jitter — 色温/光源偏移",
}


# ── 工具 ──────────────────────────────────────────────
def title_bar(slide, text, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(13.333), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PRI; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C_WHITE; p.font.name = "Arial"
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle
        p2.font.size = Pt(12); p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF); p2.font.name = "Arial"


def body_text(slide, text, left, top, width, height, size=14, bold=False, color=None, name="Arial"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color or C_DARK; p.font.name = name
    return tf


def bullets(slide, items, left, top, width, height, size=14, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.name = "Arial"
        p.font.color.rgb = color or C_DARK; p.space_after = Pt(4)
    return tf


def add_table(slide, headers, rows, left, top, width, height, col_widths=None, fs=10):
    nr = len(rows) + 1; nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(fs); p.font.bold = True; p.font.color.rgb = C_WHITE
            p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
        c.fill.solid(); c.fill.fore_color.rgb = C_PRI
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j); c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(fs - 1); p.font.color.rgb = C_DARK
                p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = C_LIGHT
    return ts


def add_image_safe(slide, fname, left, top, width, height=None):
    path = os.path.join(IMGDIR, fname)
    if os.path.exists(path):
        if height:
            return slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            return slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = f"[{fname}]"; p.font.size = Pt(10); p.font.color.rgb = C_MUTED
    return tb


def page_num(slide, n, total):
    tb = slide.shapes.add_textbox(Inches(12.2), Inches(7.15), Inches(1), Inches(0.25))
    p = tb.text_frame.paragraphs[0]; p.text = f"{n}/{total}"
    p.font.size = Pt(9); p.font.color.rgb = C_MUTED; p.font.name = "Arial"; p.alignment = PP_ALIGN.RIGHT


def set_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def accent_line(slide, top):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top),
                               Inches(5), Inches(0.03))
    s.fill.solid(); s.fill.fore_color.rgb = C_ACC; s.line.fill.background()


# ── 幻灯片 ────────────────────────────────────────────

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, C_PRI)
    body_text(s, "sys_G — 合成畴区图像生成器", 1.0, 1.5, 11, 1.2, 40, True, C_WHITE)
    accent_line(s, 2.9)
    body_text(s, "参数手册与效果展示    汇报", 1.0, 3.1, 11, 0.6, 22, color=RGBColor(0xCC, 0xDD, 0xFF))
    body_text(s, "HexagonGenerator + MicroscopyAugment  |  color_std=0.2  |  控制变量法  |  seed=42",
              1.0, 4.0, 11, 0.4, 13, color=RGBColor(0x99, 0xBB, 0xEE))
    body_text(s, "基准配置与 gen_multi_mag.py 的 SHARED_GEN_KWARGS 完全一致",
              1.0, 4.5, 11, 0.4, 12, color=RGBColor(0x88, 0xAA, 0xDD))


def slide_arch(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "系统架构概述", "sys_G Pipeline: Generator → Augment → Refine → Save")
    bullets(s, [
        "HexagonGenerator — 带噪背景上随机放置六边形畴区，OverlapController 控制重叠",
        "MicroscopyAugment — 模拟显微镜成像效果（亮度/对比度/Gamma/模糊/噪声/色温）",
        "refine_polygons() — 裁剪越界多边形，确保标注在图像边界内",
        "ISATSaver — 以 ISAT 格式保存图像和 JSON 多边形标注",
        "",
        "设计原则: 依赖倒置(DIP) + 单一职责(SRP) + 开闭原则(OCP)",
        "所有随机性通过 random.seed() / np.random.seed() 可完全复现",
    ], 0.8, 1.5, 8, 5, 14)
    bullets(s, ["① 生成畴区", "  ↓", "② 显微镜增强", "  ↓", "③ 边界裁剪", "  ↓", "④ ISAT 保存"],
            10.5, 1.5, 2.5, 5, 13, C_PRI)


def slide_multimag(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "多倍率视场与畴区密度总览",
              "radius = base_r × mag × image_scale × ss    |    count = base_n / mag²")
    headers = ["倍率", "mag_factor", "有效半径", "有效数量", "效果"]
    rows = [
        ["2.5×","0.25", "3–6 px",   "3200–6400", "远景：畴区极小、极高密度"],
        ["5×",  "0.5",  "7–12 px",  "800–1600",  "中远景：畴区小、高密度"],
        ["10×", "1.0",  "15–25 px", "200–400",   "参考倍率 (baseline)"],
        ["20×", "2.0",  "30–50 px", "50–100",    "近景：畴区大、中等密度"],
        ["50×", "5.0",  "75–125 px","8–16",      "高倍：畴区很大、稀疏"],
        ["100×","10.0", "150–250 px","2–4",      "特写：仅个别大畴区"],
    ]
    add_table(s, headers, rows, 0.3, 1.3, 12, 2.6, [0.8,1.2,1.5,1.5,3.2], 10)
    add_image_safe(s, "multi_mag_overview.png", 0.3, 4.2, 12.5, 2.8)
    body_text(s, "同一套基准参数仅通过 mag_factor 产生完全不同的视场效果  |  "
              "低倍率 → 畴区小密集 (统计分布)  |  高倍率 → 畴区大清晰 (形态学)",
              0.5, 7.15, 12, 0.25, 10, color=C_MUTED)


def slide_param_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "生成器参数总览", "与 gen_multi_mag.py SHARED_GEN_KWARGS 一致  |  color_std=0.2")
    headers = ["分组", "参数", "类型", "默认值", "含义"]
    rows = [
        ["颜色/纹理","color_mean","tuple","(163,138,127)","畴区 RGB 均值"],
        ["","bg_mean","tuple","(153,115,85)","背景 RGB 均值"],
        ["","color_std","float","0.2 (固定)","畴区间颜色差异"],
        ["","texture_std","float","2","畴区内部纹理噪声"],
        ["","bg_noise_std","float","2","背景像素噪声"],
        ["几何/尺寸","base_r_range","tuple","(15,25)","基础半径范围"],
        ["","base_num_range","tuple","(200,400)","基础数量范围"],
        ["","size_std","float","25","大小分布标准差"],
        ["","mag_factor","float","1.0","缩放倍率（联动数量）"],
        ["","shape_jitter","float","0.1","顶点扰动比例"],
        ["","orientation_std","float","0","畴区取向标准差（度）"],
        ["","image_scale","float","auto","画布缩放因子 (=W/1024)"],
        ["边缘毛刺","edge_burr_amp","float","0.08","毛刺振幅"],
        ["","edge_burr_subdiv","int","4","边细分点数"],
        ["重叠控制","max_overlap_ratio","float","0.3","最大重叠比例"],
        ["","max_overlap_count","int","3","最大覆盖层数"],
        ["","contain_threshold","float","0.85","包含判定阈值"],
        ["渲染","supersample_ratio","int","1","超采样倍率"],
    ]
    add_table(s, headers, rows, 0.3, 1.2, 12.5, 5.8, [1.2,2.0,1.0,2.0,4.0], 9)
    body_text(s, "18 个可调参数，5 个功能组  |  配合 7 组增强参数 = 完整的合成数据控制能力",
              0.5, 7.15, 10, 0.25, 11, color=C_MUTED)


def slide_aug_param(prs, aug_key, pg, total):
    """单个增强器参数幻灯片。"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    display_name = AUG_PARAM_NAMES.get(aug_key, aug_key)
    title_bar(s, display_name, "MicroscopyAugment  |  各项关闭，逐一开启测试")

    expl = AUG_EXPLANATIONS.get(aug_key, "")
    body_text(s, expl, 0.5, 1.25, 12, 1.0, 13)

    fname = AUG_PARAM_IMG_MAP.get(aug_key, "")
    add_image_safe(s, fname, 0.3, 2.3, 12.5, 4.6)

    page_num(s, pg[0], total)
    pg[0] += 1


def slide_interaction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "参数交互关系", "核心缩放公式  |  重叠协同  |  毛刺+形状组合")

    body_text(s, "半径与数量的联合缩放", 0.5, 1.3, 5.5, 0.4, 18, True, C_PRI)
    body_text(s,
        "effective_radius = base_r × mag_factor\n"
        "                    × image_scale × supersample_ratio\n\n"
        "effective_count  = base_num / mag_factor²\n\n"
        "• mag_factor ↑ → 畴区变大、数量减少 (近景)\n"
        "• mag_factor ↓ → 畴区变小、数量增加 (远景)\n"
        "• image_scale → 仅影响半径，不影响数量",
        0.5, 1.8, 5.5, 3.2, 12)

    body_text(s, "重叠参数协同", 6.8, 1.3, 5.5, 0.4, 18, True, C_PRI)
    headers = ["场景", "overlap_ratio", "overlap_count", "contain_threshold"]
    rows = [
        ["完全分离","0.0","1","0.5"],
        ["轻微接触","0.15-0.25","2","0.85"],
        ["适度重叠(默认)","0.3","3","0.85"],
        ["重度重叠","0.6+","None","0.95"],
    ]
    add_table(s, headers, rows, 6.8, 1.8, 5.5, 1.6, [1.5,1.5,1.3,1.2], 10)

    body_text(s, "边缘毛刺 + 形状扰动", 6.8, 3.8, 5.5, 0.4, 18, True, C_PRI)
    body_text(s,
        "shape_jitter: 六边形整体形状 (低频)\n"
        "edge_burr_*:  每条边高频毛刺细节\n\n"
        "当前默认值:\n"
        "  shape_jitter = 0.1\n"
        "  edge_burr_amplitude = 0.08\n"
        "  edge_burr_subdivisions = 4\n"
        "→ 产生自然的晶粒边缘效果",
        6.8, 4.3, 5.5, 2.5, 12)


def slide_aug_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "增强器参数总览 — MicroscopyAugment",
              "与 gen_multi_mag.py SHARED_AUG_KWARGS 一致  |  模拟显微镜成像退化")
    headers = ["参数组", "字段", "默认值", "含义"]
    rows = [
        ["亮度","brightness_range / prob","(0.8,1) / 0.8","仅向暗调变化"],
        ["对比度","contrast_range / prob","(0.4,1) / 0.8","仅向低对比度变化"],
        ["Gamma","gamma_range / prob","(0.7,1.3) / 0.8","非线性亮度映射"],
        ["模糊","blur_range / prob","(0.5,1) / 1.0","高斯模糊 (始终触发)"],
        ["色温","color_jitter_range / prob","(-12,12) / 0.8","RGB 偏移 (最先施加)"],
        ["椒盐噪声","sp_noise_prob","0.0 (关闭)","传感器坏点模拟"],
        ["旋转","rotate_prob","0.0 (关闭)","图像旋转增强"],
    ]
    add_table(s, headers, rows, 0.5, 1.3, 12, 3.0, [1.2,3.0,2.5,3.5], 11)
    bullets(s, [
        "每个增强项有独立触发概率 prob，设为 0 可完全禁用",
        "color_jitter (色温) 在所有其他增强之前施加，模拟光源变化",
        "增强器在生成器之后、保存之前执行；旋转后多边形坐标自动跟随",
    ], 0.8, 4.7, 11, 2.0, 13)


def slide_workflow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "调参工作流建议", "从真实图像出发 → 控制变量调参 → 批量前验证")

    body_text(s, "Step 1: 从真实图像出发", 0.5, 1.3, 12, 0.4, 18, True, C_PRI)
    bullets(s, [
        "① 确定倍率 → 从 MAG_FACTORS 选择 mag_factor (2.5x-100x)",
        "② 匹配颜色 → 采样畴区/背景平均 RGB → color_mean / bg_mean",
        "③ 估计尺寸 → 测量畴区直径 (px) → 反推 base_r_range",
        "④ 估计数量 → 统计畴区数量 → 反推 base_num_range",
        "⑤ 微调纹理 → texture_std, bg_noise_std, edge_burr_amplitude",
    ], 0.5, 1.8, 12, 2.0, 12)

    body_text(s, "Step 2: 控制变量调参", 0.5, 3.9, 12, 0.4, 18, True, C_PRI)
    bullets(s, [
        "python -m om_domain.param_showcase --size 1024       # 生成全部对比图",
        "python -m om_domain.param_showcase --multi-mag-only  # 仅多倍率总览",
        "python -m om_domain.param_showcase --gen-only        # 仅生成器参数",
        "python -m om_domain.param_showcase --aug-only        # 仅增强器参数",
    ], 0.5, 4.4, 7.5, 1.5, 12, C_PRI)

    body_text(s, "Step 3: 批量生成前验证", 0.5, 6.0, 12, 0.4, 18, True, C_PRI)
    body_text(s, "pipeline.run('./output/test', n=5) → 人工确认 → gen_multi_mag.py 批量产出",
              0.5, 6.5, 12, 0.4, 12)


def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, C_PRI)
    body_text(s, "总结", 0.8, 0.8, 4, 0.6, 36, True, C_WHITE)
    accent_line(s, 1.5)
    bullets(s, [
        "sys_G = HexagonGenerator + MicroscopyAugment + ISATSaver",
        "",
        "18 个生成器参数 × 7 组增强参数 = 完整的合成数据控制能力",
        "",
        "核心公式: radius ∝ mag  |  count ∝ 1/mag²  |  image_scale = W/1024",
        "",
        "设计原则: 控制变量法验证 + 参数正交化 + 完全可复现",
        "",
        "多倍率生成: gen_multi_mag.py 一键产出 2.5x-100x 六种倍率数据集",
        "",
        "可视化工具: param_showcase.py 展示全部参数效果",
    ], 0.8, 1.8, 11, 5.3, 15, C_WHITE)


def slide_config(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "附录: 基准配置 (与 gen_multi_mag.py 一致)", "color_std=0.2")
    code = (
        "BASELINE_CONFIG = dict(\n"
        "    color_mean=(163,138,127), bg_mean=(153,115,85),\n"
        "    color_std=0.2, texture_std=2, bg_noise_std=2,\n"
        "    base_r_range=(15,25), base_num_range=(200,400),\n"
        "    size_std=25, mag_factor=1.0, shape_jitter=0.1,\n"
        "    orientation_std=0, image_scale=None,  # auto = W/1024\n"
        "    edge_burr_amplitude=0.08, edge_burr_subdivisions=4,\n"
        "    max_overlap_ratio=0.3, max_overlap_count=3,\n"
        "    contain_threshold=0.85, min_area_factor=5,\n"
        "    supersample_ratio=1,\n"
        ")\n\n"
        "MAG_FACTORS = {\n"
        "    '2.5x':0.25, '5x':0.5, '10x':1.0,\n"
        "    '20x':2.0, '50x':5.0, '100x':10.0\n"
        "}"
    )
    body_text(s, code, 0.5, 1.3, 12, 5.8, 11, name="Courier New")


# ── 主入口 ────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="生成 sys_G 参数汇报 PPT")
    parser.add_argument("-o", "--output", default="./output/sys_G_report.pptx")
    parser.add_argument("--size", default="16:9", choices=["16:9", "4:3"])
    args = parser.parse_args()

    out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), args.output)
    os.makedirs(os.path.dirname(out), exist_ok=True)

    prs = Presentation()
    if args.size == "16:9":
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)

    # 生成器参数列表
    gen_params = [
        ("颜色纹理", "color_std"),
        ("颜色纹理", "texture_std"),
        ("颜色纹理", "bg_noise_std"),
        ("颜色纹理", "color_mean"),
        ("颜色纹理", "bg_mean"),
        ("几何尺寸", "base_r_range"),
        ("几何尺寸", "shape_jitter"),
        ("几何尺寸", "orientation_std"),
        ("几何尺寸", "size_std"),
        ("几何尺寸", "base_num_range"),
        ("边缘毛刺", "edge_burr_amplitude"),
        ("边缘毛刺", "edge_burr_subdivisions"),
        ("重叠控制", "max_overlap_ratio"),
        ("重叠控制", "max_overlap_count"),
        ("重叠控制", "contain_threshold"),
        ("渲染", "supersample_ratio"),
    ]
    aug_keys = ["Brightness", "Contrast", "Gamma", "Gaussian", "Salt", "Color"]

    # 计算总页数
    total = 4 + len(gen_params) + 1 + len(aug_keys) + 3 + 1  # cover+arch+multi+table + gen* + aug_overview + aug* + inter+workflow+summary + config
    pg = [1]  # mutable counter

    # 封面
    slide_title(prs)

    # 架构
    slide_arch(prs); page_num(prs.slides[-1], pg[0], total); pg[0] += 1

    # 多倍率
    slide_multimag(prs); page_num(prs.slides[-1], pg[0], total); pg[0] += 1

    # 参数总览表
    slide_param_table(prs); page_num(prs.slides[-1], pg[0], total); pg[0] += 1

    # 每个生成器参数
    for cat, pname in gen_params:
        slide_gen_param_direct(prs, cat, pname, pg, total)

    # 增强器总览
    slide_aug_overview(prs); page_num(prs.slides[-1], pg[0], total); pg[0] += 1

    # 每个增强器参数
    for ak in aug_keys:
        slide_aug_param(prs, ak, pg, total)

    # 交互关系
    slide_interaction(prs); page_num(prs.slides[-1], pg[0], total); pg[0] += 1

    # 工作流
    slide_workflow(prs); page_num(prs.slides[-1], pg[0], total); pg[0] += 1

    # 总结
    slide_summary(prs)

    # 附录
    slide_config(prs); page_num(prs.slides[-1], pg[0], total); pg[0] += 1

    prs.save(out)
    print(f"PPT 已保存: {os.path.abspath(out)}")
    print(f"共 {total} 张幻灯片")


def slide_gen_param_direct(prs, cat, pname, pg, total):
    """直接构建生成器参数幻灯片。"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, pname, f"分组: {cat}  |  color_std=0.2  |  控制变量法 (seed=42)")

    expl = GEN_EXPLANATIONS.get(pname, "")
    body_text(s, expl, 0.5, 1.2, 12, 1.0, 13)

    fname = GEN_PARAM_IMG_MAP.get(pname, "")
    add_image_safe(s, fname, 0.2, 2.3, 12.8, 4.9)

    page_num(s, pg[0], total)
    pg[0] += 1


if __name__ == "__main__":
    main()
